# -*- coding: utf-8 -*-
"""生成答辩 Keynote(.pptx,Keynote 可直接打开)。清华紫风格,中文 PingFang SC。
文字原生可编辑;架构/闭环/四层/甘特 4 张图用 beamer 渲染图嵌入。"""
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE="/home/lizhaoxi/.kitty/workspaces/s_9205bdbc/kaiti"
IMG=BASE+"/keynote/img"
PURPLE=RGBColor(0x74,0x34,0x81); BLUE=RGBColor(0x00,0x4E,0x98)
DARK=RGBColor(0x22,0x22,0x22); GRAY=RGBColor(0x70,0x70,0x70); WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="PingFang SC"

# render logo to png
d=fitz.open(BASE+"/slides/thu_logo.pdf"); d[0].get_pixmap(matrix=fitz.Matrix(4,4)).save(IMG+"/logo.png"); d.close()

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def setrun(run,size,bold=False,color=DARK,name=FONT):
    f=run.font; f.size=Pt(size); f.bold=bold; f.color.rgb=color; f.name=name
    rPr=run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(tag))
        if e is None: e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',name)

def blank():
    return prs.slides.add_slide(BLANK)

def band(slide,title,sub=None):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(1.02))
    bar.fill.solid(); bar.fill.fore_color.rgb=PURPLE; bar.line.fill.background()
    bar.shadow.inherit=False
    tb=slide.shapes.add_textbox(Inches(0.55),Inches(0.12),SW-Inches(1.1),Inches(0.8)).text_frame
    tb.word_wrap=True; p=tb.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    setrun(p.add_run(),26,True,WHITE); p.runs[0].text=title
    tb.vertical_anchor=MSO_ANCHOR.MIDDLE

def body(slide,top=1.3,left=0.7,width=None,height=None):
    width=width or SW-Inches(1.4); height=height or SH-Inches(1.7)
    box=slide.shapes.add_textbox(Inches(left),Inches(top),width if isinstance(width,Emu) else Inches(width),Inches(height))
    tf=box.text_frame; tf.word_wrap=True; return tf

def bullets(tf,items,base=18):
    # items: (level, text, bold, color)
    first=True
    for it in items:
        lvl,txt=it[0],it[1]; bold=it[2] if len(it)>2 else False; col=it[3] if len(it)>3 else DARK
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level=lvl; p.space_after=Pt(6); p.line_spacing=1.12
        # bullet char
        prefix={0:"▸ ",1:"• ",2:"– "}.get(lvl,"")
        r=p.add_run(); r.text=prefix+txt; setrun(r,base-lvl*1.5,bold,col)

def figure_slide(title,img,note=None):
    s=blank(); band(s,title)
    im=IMG+f"/{img}.png"
    from PIL import Image
    try:
        iw,ih=Image.open(im).size
    except Exception:
        px=fitz.open(); iw,ih=1600,800
    avail_w=Inches(12.2); avail_h=Inches(5.4 if note else 5.9)
    scale=min(avail_w/iw, avail_h/ih)
    w=int(iw*scale); h=int(ih*scale)
    left=int((SW-w)/2); top=Inches(1.25)
    s.shapes.add_picture(im,left,top,width=w,height=h)
    if note:
        tf=body(s,top=6.6,left=0.7,height=0.7); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=note; setrun(r,13,False,GRAY)
    return s

# ===== S1 封面 =====
s=blank()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(2.55)); bar.fill.solid(); bar.fill.fore_color.rgb=PURPLE; bar.line.fill.background(); bar.shadow.inherit=False
tf=s.shapes.add_textbox(Inches(0.8),Inches(0.55),SW-Inches(1.6),Inches(1.9)).text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text="基于大模型编码智能体的\n自进化 eBPF 网络入侵防御系统"; setrun(r,30,True,WHITE)
p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; r=p2.add_run(); r.text="硕士学位论文开题报告"; setrun(r,18,False,RGBColor(0xEE,0xE0,0xF2))
try: s.shapes.add_picture(IMG+"/logo.png",int((SW-Inches(1.2))/2),Inches(2.95),height=Inches(1.2))
except Exception: pass
tf2=s.shapes.add_textbox(Inches(0.8),Inches(4.5),SW-Inches(1.6),Inches(2.2)).text_frame
for t,sz,bd in [("研究生:李兆曦      指导教师:徐恪 教授",18,True),("清华大学 计算机科学与技术系",16,False),("2026 年 6 月",16,False)]:
    p=tf2.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_after=Pt(8); r=p.add_run(); r.text=t; setrun(r,sz,bd,PURPLE if bd else GRAY)

# ===== S2 目录 =====
s=blank(); band(s,"目录")
tf=body(s); items=[(0,"一、研究背景与动机"),(0,"二、研究问题与挑战"),(0,"三、相关工作"),(0,"四、研究内容与方法"),
 (0,"五、技术路线"),(0,"六、创新点"),(0,"七、工作计划"),(0,"八、预期成果")]
bullets(tf,items,base=22)

# ===== S3 背景① =====
s=blank(); band(s,"研究背景:跨层歧义持续催生 off-path 协议攻击")
bullets(body(s),[
 (0,"TCP/IP 跨层交互存在语义歧义,长期催生 off-path(链外)攻击",True),
 (1,"攻击者不在路径上、不窃听,仅凭侧信道即可推断序列号 → 劫持/重置/污染"),
 (1,"例:challenge-ACK(CVE-2016-5696)、混合 IPID、PMTUD 分片、NAT/Wi-Fi 劫持"),
 (0,"近期 CCS'25:PMTUD 打破 IP 共享场景下的 TCP 连接隔离",True),
 (1,"平均 220 秒、约 70% 成功率推断序列号;50 个真实网络中 38 个脆弱"),
 (0,"核心判断:攻击面植根于设计层面的跨层歧义,不断被重新发掘",True,PURPLE),
])

# ===== S4 背景② =====
s=blank(); band(s,"规律:防御机制本身在制造新的攻击面")
bullets(body(s),[
 (0,"很多新攻击面恰由防御机制自身引入——“补丁”反成下一处侧信道",True),
 (1,"challenge-ACK 限速 → 侧信道;源端口随机化 → 被反推"),
 (0,"启示:攻击与防御同源、相互转化",True,PURPLE),
 (1,"防御必须持续演化,一次性静态修补难以为继"),
])

# ===== S5 背景③ =====
s=blank(); band(s,"现有防御的三大结构性缺陷")
bullets(body(s),[
 (0,"零散:各攻击面分别打补丁,缺乏统一框架",True),
 (0,"滞后:先有攻击曝光、再被动响应,披露到上线窗口长",True),
 (0,"依赖人工编写内核代码",True),
 (1,"周期长、易出错;经验留于个人,难以沉淀复用"),
 (0,"症结:防御的生产方式仍是人工、低速、难以累积的",True,PURPLE),
])

# ===== S6 背景④ =====
s=blank(); band(s,"两项技术趋于成熟,使自动化防御成为可能")
bullets(body(s),[
 (0,"eBPF 内核可编程",True),
 (1,"热加载(不改内核源码、不重启)、贴近协议栈(XDP/TC 线速)"),
 (1,"内置 verifier:加载前做形式化安全验证(有界、内存安全、可终止)"),
 (0,"大模型编码智能体",True),
 (1,"能自主编写/修改/调试系统级代码,具工具调用与反思迭代"),
 (1,"代表:Claude Code、Cursor、Codex、Devin、SWE-agent"),
 (0,"契机:二者结合,使“由智能体持续维护内核防御”可行",True,PURPLE),
])

# ===== S7 核心问题 =====
s=blank(); band(s,"核心研究问题")
tf=body(s,top=1.7)
p=tf.paragraphs[0]; r=p.add_run(); r.text="能否让大模型编码智能体,自主维护一套运行在内核的 eBPF 防御库,\n使其持续进化,以应对不断涌现的协议攻击?"; setrun(r,24,True,PURPLE); p.line_spacing=1.3
bullets(body(s,top=4.0),[
 (0,"目标不是“一次性合成单个防御程序”,而是长期自主运转的防御系统"),
 (0,"核心:自主性(探索—合成—验证—部署闭环)与进化性(能力单调累积、不回退)"),
])

# ===== S8 挑战 =====
s=blank(); band(s,"关键挑战(工程上)")
bullets(body(s),[
 (0,"可靠性:智能体写的内核代码若有缺陷可断网/误杀 → 先验证安全、再上线,出错能回滚",True),
 (0,"生成质量:智能体写的 eBPF 要真能编过、过 verifier、且确实拦住攻击而不误杀",True),
 (0,"评测:怎么证明真有用——见过的、没见过的攻击都能防,且开销可接受",True),
])

# ===== S9 相关工作① =====
s=blank(); band(s,"相关工作(一):协议攻防与 eBPF 运行时防御")
bullets(body(s),[
 (0,"协议攻击与检测",True),
 (1,"攻击:off-path 系列(challenge-ACK、IPID、PMTUD 等);检测:SCENT、SCAD"),
 (1,"现状:防御零散、滞后,缺统一且持续演化的框架"),
 (0,"eBPF 运行时防御",True),
 (1,"系统:Katran、Cilium、Falco、Tetragon——能力强,但策略均人工编写"),
 (1,"验证:PREVAIL、K2 等围绕 verifier 的程序分析/优化"),
])

# ===== S10 相关工作② =====
s=blank(); band(s,"相关工作(二):大模型用于安全 / 编码智能体")
bullets(body(s),[
 (0,"大模型 × 安全",True),
 (1,"进攻:PentestGPT、Big Sleep"),
 (1,"合成 eBPF:Kgent / KEN——最直接对比基线,但一次性合成、不持续维护"),
 (0,"编码智能体(工程载体)",True),
 (1,"Claude Code/Cursor/Codex/Devin/SWE-agent:能自主写改调代码,但通用、不带内核/eBPF 领域工具——本文做专用 harness"),
 (1,"迭代思路可借鉴 FunSearch、AlphaEvolve(让模型反复改代码、用测试反馈迭代)"),
])

# ===== S11 定位对比(表) =====
s=blank(); band(s,"定位对比:本文与代表性工作")
rows=[("对比对象","已有范式","本文范式"),
 ("Kgent / KEN","一次性合成 eBPF","长期自主维护的防御库"),
 ("Tetragon / Falco","人工编写策略","自进化地生成与更新策略"),
 ("小模型(权重学习)","知识沉淀为权重","知识沉淀为可验证代码")]
tb=s.shapes.add_table(4,3,Inches(0.9),Inches(1.7),Inches(11.5),Inches(3.2)).table
tb.columns[0].width=Inches(3.3); tb.columns[1].width=Inches(4.1); tb.columns[2].width=Inches(4.1)
for r in range(4):
    for c in range(3):
        cell=tb.cell(r,c); cell.text=""
        p=cell.text_frame.paragraphs[0]; run=p.add_run(); run.text=rows[r][c]
        setrun(run,16,(r==0),WHITE if r==0 else DARK)
        cell.fill.solid(); cell.fill.fore_color.rgb=PURPLE if r==0 else (RGBColor(0xF2,0xEC,0xF5) if c==2 else WHITE)
tf=body(s,top=5.2,height=1.0); p=tf.paragraphs[0]; r=p.add_run()
r.text="共同差异:由“一次合成 / 人工 / 权重” → 持续、自主、可验证的代码演化"; setrun(r,16,True,PURPLE)

# ===== S12 核心思路 =====
s=blank(); band(s,"总体思路:让智能体持续维护内核 eBPF 防御")
tf=body(s,top=1.5,height=1.6); p=tf.paragraphs[0]; p.line_spacing=1.25
r=p.add_run(); r.text="把“发现攻击 → 写/改 eBPF → 验证安全 → 加载 → 看效果 → 再改”做成一个自动闭环;防御以 eBPF 代码的形式不断积累在“防御库”里(项目内部也称“防御演化”)。"; setrun(r,19,False,DARK)
bullets(body(s,top=3.6),[
 (0,"为什么用“代码”而不是训练一个模型",True),
 (1,"改一条、加一条,可积累、不遗忘;每条防御就是一段代码,看得懂、可审计"),
 (1,"跑在 XDP,线速;出错能回滚,verifier 保证不把内核搞挂"),
])

# ===== S13 研究内容(三项) =====
s=blank(); band(s,"研究内容(三项,落在系统上)")
bullets(body(s),[
 (0,"把系统搭起来:智能体 + eBPF 工具链闭环",True),
 (1,"在 PacketScope/Guarder 上,做一套让智能体能“生成 eBPF→编译→过 verifier→沙箱测试→加载→读遥测”的工具与闭环"),
 (0,"让智能体写出“能用且安全”的 eBPF",True),
 (1,"解决生成质量(编得过、过 verifier、不误杀)与安全(出错自动回滚)——系统能跑起来的关键工程"),
 (0,"在真实协议攻击上做防御 + 评测",True),
 (1,"用 PMTUD/NAT/IPID 等真实攻击搭靶场,测检出率、误报率、性能开销,与人工规则/现有工具对比"),
])

# ===== S14 架构图 =====
figure_slide("系统分层架构","arch","控制面(智能体)生成候选 → 可信验证面把关 → 内核数据面执行;反馈回灌驱动持续迭代,有效防御沉淀入知识库。")
# ===== S15 闭环 =====
figure_slide("技术路线:防御演化闭环(五阶段)","loop","每轮迭代将新发现的攻击面转化为一段经 verifier 背书的 eBPF 防御,并沉淀入防御库。")
# ===== S16 四层 =====
figure_slide("可信保障:四层安全约束 + 自动回滚","layer","这四层都是现成基础设施;创新在于把它们接入智能体闭环:报错驱动自修正、沙箱把关、异常自动回滚。")

# ===== S17 工程基线 =====
s=blank(); band(s,"工程基线:已有开源系统积累")
bullets(body(s),[
 (0,"PacketScope / Guarder",True),
 (1,"XDP 数据面(网卡驱动层线速处理报文)"),
 (1,"bpf2go 一键编译;规则热加载(运行时动态更新,无需重启)"),
 (1,"已具备 Agent 可编程 eBPF 原型,配套 312 项测试"),
 (0,"已具备数据面、编译链与热加载基础设施,闭环可在现有底座上快速搭建",True,PURPLE),
])

# ===== S18 评测 =====
s=blank(); band(s,"怎么评测")
bullets(body(s),[
 (0,"思路:用真实协议攻击搭靶场,让系统自动生成防御;并对“没见过的新攻击”单独测一遍",True),
 (0,"看四个指标",True),
 (1,"防住没(检出率)、误杀多少(误报率)"),
 (1,"开销(XDP 吞吐 / 时延)、要不要人工(自动化程度)"),
 (0,"跟谁比",True),
 (1,"人工写的规则(现状)、Kgent/KEN(一次合成)、Tetragon/Falco(人工策略)"),
])

# ===== S19 创新点 =====
s=blank(); band(s,"创新点")
bullets(body(s,top=1.6),[
 (0,"一套能跑起来的系统(而非设想):在 PacketScope/Guarder 上落地,让智能体自动维护内核 eBPF 防御,可演示可复现",True),
 (0,"智能体写内核代码的可靠流水线:生成→编译→verifier→沙箱→加载→回滚,保证“AI 写的 eBPF”能用、且不把内核搞挂",True),
 (0,"真实攻击的防御 + 评测靶场:用 PMTUD/NAT/IPID 等真实攻击,给出检出率/误报/开销,并与人工规则、现有工具对比",True),
],base=19)

# ===== S20 工作计划 =====
s=blank(); band(s,"工作计划(2026.06 – 2027.06)")
bullets(body(s,top=1.3,height=2.5),[
 (0,"2026.07–09:防御演化框架原型(形式化 + 系统骨架)"),
 (0,"2026.10–12:闭环原型 + 初步实验(探索/合成/验证/部署/档案)"),
 (0,"2027.01–03:评测 + 对比实验 + 中期检查"),
 (0,"2027.04–05:论文撰写、预答辩、送审　|　2027.06:答辩"),
],base=17)
from PIL import Image as _I
iw,ih=_I.open(IMG+"/gantt.png").size; scale=min(Inches(11.0)/iw,Inches(2.2)/ih)
s.shapes.add_picture(IMG+"/gantt.png",int((SW-int(iw*scale))/2),Inches(4.6),width=int(iw*scale),height=int(ih*scale))

# ===== S21 预期成果 =====
s=blank(); band(s,"预期成果")
bullets(body(s),[
 (0,"学位论文:完成硕士学位论文 1 篇",True),
 (0,"开源系统:自进化 eBPF 防御系统、防御演化专用 harness、攻防评测基准",True),
 (0,"社区贡献:对发现的协议/实现缺陷,向 IETF / Linux 负责任披露",True),
])

# ===== S22 致谢 =====
s=blank()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); bar.fill.solid(); bar.fill.fore_color.rgb=PURPLE; bar.line.fill.background(); bar.shadow.inherit=False
tf=s.shapes.add_textbox(Inches(1),Inches(2.7),SW-Inches(2),Inches(2)).text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text="敬请各位老师批评指正"; setrun(r,34,True,WHITE)
p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; p2.space_before=Pt(20); r=p2.add_run(); r.text="谢谢！"; setrun(r,24,False,RGBColor(0xEE,0xE0,0xF2))

OUT=BASE+"/keynote/答辩Keynote_李兆曦.pptx"
prs.save(OUT)
print("SAVED",OUT,"slides:",len(prs.slides._sldIdLst))
